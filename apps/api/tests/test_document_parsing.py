"""Upload file -> Markdown: dinh tuyen, cache va chuyen doi bang MarkItDown (cuc bo)."""

from __future__ import annotations

import zipfile
from pathlib import Path
from types import SimpleNamespace
from typing import Any

import pytest

from sag_api.core.config import Settings
from sag_api.parsing import service
from sag_api.parsing.service import PreparedDocument
from sag_api.sag.dto import ProcessCheckpoint, ProcessOutcome


def _settings(**overrides: Any) -> Settings:
    return Settings(
        _env_file=None,
        data_dir="/tmp/sag-test-engine",
        upload_dir="/tmp/sag-test-uploads",
        **overrides,
    )


@pytest.mark.asyncio
async def test_parser_routes_markdown_and_markitdown_with_cache(tmp_path, monkeypatch):
    markdown = tmp_path / "already.md"
    markdown.write_text("# Already\n", encoding="utf-8")
    direct = await service.prepare_document(str(markdown), _settings())
    assert direct.path == str(markdown)
    assert direct.provider == "original"

    source = tmp_path / "notes.docx"
    source.write_bytes(b"fake-office")
    calls: list[str] = []

    def convert(path: str) -> str:
        calls.append(path)
        return "# Converted\n\nhello"

    monkeypatch.setattr(service, "_markitdown_sync", convert)
    first = await service.prepare_document(str(source), _settings())
    second = await service.prepare_document(str(source), _settings())

    assert first.provider == "markitdown" and first.path.endswith(".parsed.markitdown.md")
    assert Path(first.path).read_text(encoding="utf-8") == "# Converted\n\nhello\n"
    assert second.cached is True and second.path == first.path
    assert calls == [str(source)]


@pytest.mark.asyncio
async def test_legacy_gb18030_text_is_normalized_without_markitdown(tmp_path, monkeypatch):
    source = tmp_path / "\u9a86\u9a7c\u7965\u5b50.txt"
    expected = "\u300a\u9a86\u9a7c\u7965\u5b50\u300b\r\n\u4f5c\u8005\uff1a\u8001\u820d\r\n\u6b63\u6587\u53ea\u6709\u4e00\u4e2a\u635f\u574f\u5b57\u8282\uff1a"
    source.write_bytes(expected.encode("gb18030") + b"\xff")
    stale_cache = Path(f"{source}.parsed.markitdown.md")
    stale_cache.write_text("None\n", encoding="utf-8")

    def should_not_run(_path: str) -> str:
        raise AssertionError("plain text should use Muse's text decoder")

    monkeypatch.setattr(service, "_markitdown_sync", should_not_run)
    parsed = await service.prepare_document(str(source), _settings())
    markdown = Path(parsed.path).read_text(encoding="utf-8")

    assert parsed.cached is False
    assert parsed.provider == "markitdown"
    assert markdown.startswith(expected.replace("\r\n", "\n"))
    assert markdown.count("�") == 1
    assert markdown != "None\n"


@pytest.mark.asyncio
async def test_markitdown_none_sentinel_is_rejected(tmp_path, monkeypatch):
    source = tmp_path / "broken.docx"
    source.write_bytes(b"fake-office")
    monkeypatch.setattr(service, "_markitdown_sync", lambda _path: "None")

    with pytest.raises(Exception, match="kh\u00f4ng l\u1ea5y \u0111\u01b0\u1ee3c n\u1ed9i dung h\u1ee3p l\u1ec7"):
        await service.prepare_document(str(source), _settings())

    assert not Path(f"{source}.parsed.markitdown.md").exists()


@pytest.mark.asyncio
async def test_document_job_sends_parsed_markdown_to_engine(monkeypatch):
    from sag_api.db.models import Document, Source
    from sag_api.jobs import tasks

    document = SimpleNamespace(
        id="doc-1",
        source_id="source-1",
        filename="original.pdf",
        storage_path="/uploads/original.pdf",
        status=None,
        error=None,
        chunk_count=0,
        event_count=0,
        progress=0,
        token_usage=0,
        sag_source_id=None,
    )
    source = SimpleNamespace(
        id="source-1",
        sag_source_config_id="sag-source-1",
        chunk_count=0,
        event_count=0,
    )
    job = SimpleNamespace(id="job-1", document_id="doc-1", progress=0.0, payload={})

    class FakeSession:
        async def get(self, model, _id):
            return document if model is Document else source if model is Source else None

        async def commit(self):
            pass

        async def execute(self, _statement):
            pass

        async def refresh(self, _instance, attribute_names=None):
            pass

    prepared_calls: list[str] = []

    async def fake_prepare(path, settings, *, state=None, on_state=None):
        prepared_calls.append(path)
        return PreparedDocument("/uploads/original.pdf.parsed.markitdown.md", "markitdown")

    class FakeEngineManager:
        seen_path = ""

        async def process_document(
            self,
            source_config_id,
            path,
            *,
            source,
            on_stage,
            checkpoint,
            on_checkpoint,
            should_pause,
            max_concurrency,
            document_title,
        ):
            self.seen_path = path
            assert max_concurrency == tasks.settings.document_extract_concurrency
            assert document_title == "original"
            await on_stage("loading")
            await on_checkpoint(
                ProcessCheckpoint(
                    source_id="engine-doc",
                    chunk_ids=["chunk-1", "chunk-2"],
                    processed_chunk_ids=["chunk-1"],
                    event_count=1,
                    event_ids=["event-1"],
                    token_usage=1234,
                )
            )
            await on_stage("extracting")
            return ProcessOutcome(
                source_id="engine-doc",
                chunk_count=2,
                event_count=1,
                chunk_ids=["chunk-1", "chunk-2"],
                processed_chunk_ids=["chunk-1", "chunk-2"],
                token_usage=2468,
            )

    monkeypatch.setattr(tasks, "prepare_document", fake_prepare)
    engine = FakeEngineManager()
    await tasks.process_document(FakeSession(), job, engine_manager=engine)

    assert prepared_calls == ["/uploads/original.pdf"]
    assert engine.seen_path.endswith(".md")
    assert document.status.value == "ready"
    assert document.chunk_count == 2 and document.event_count == 1
    assert document.progress == 100 and document.token_usage == 2468


def _simple_pdf(text: str) -> bytes:
    """Build a minimal PDF with an extractable text layer, so the test needs no PDF writing library."""
    stream = f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, obj in enumerate(objects, 1):
        offsets.append(len(output))
        output.extend(f"{number} 0 obj\n".encode() + obj + b"\nendobj\n")
    xref = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode())
    output.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode()
    )
    return bytes(output)


def _simple_docx(path: Path, text: str) -> None:
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
              <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
              <Default Extension="xml" ContentType="application/xml"/>
              <Override PartName="/word/document.xml"
                ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
            </Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
              <Relationship Id="rId1"
                Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument"
                Target="word/document.xml"/>
            </Relationships>""",
        )
        archive.writestr(
            "word/document.xml",
            f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
            <w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
              <w:body><w:p><w:r><w:t>{text}</w:t></w:r></w:p></w:body>
            </w:document>""",
        )


def test_real_markitdown_converts_pdf_and_office_files(tmp_path):
    """Dependency install smoke test: the core formats really do produce Markdown the engine can ingest."""
    from openpyxl import Workbook
    from pptx import Presentation

    from sag_api.parsing.service import _markitdown_sync

    pdf = tmp_path / "sample.pdf"
    pdf.write_bytes(_simple_pdf("Muse PDF marker"))

    docx = tmp_path / "sample.docx"
    _simple_docx(docx, "Muse DOCX marker")

    pptx = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Muse PPTX marker"
    presentation.save(pptx)

    xlsx = tmp_path / "sample.xlsx"
    workbook = Workbook()
    workbook.active["A1"] = "Muse XLSX marker"
    workbook.save(xlsx)

    assert "Muse PDF marker" in _markitdown_sync(str(pdf))
    assert "Muse DOCX marker" in _markitdown_sync(str(docx))
    assert "Muse PPTX marker" in _markitdown_sync(str(pptx))
    assert "Muse XLSX marker" in _markitdown_sync(str(xlsx))
